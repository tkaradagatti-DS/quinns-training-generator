"""
Enhanced Training Generator - Complete Production System with Modern UI
A comprehensive AI-powered training material generator with 4-phase workflow

Version: 3.2 Production (Enhanced UI + Theme Toggle + Secrets.toml)
Author: QUINNS Training Services Limited
License: Proprietary

ENHANCEMENTS:
- Dual theme support (Dark/Light mode toggle)
- Color-coded importance levels
- Improved text visibility throughout
- Visible number input controls
- Always-visible download icons
- Auto-load API key from .streamlit/secrets.toml file
"""

# ============================================================================
# IMPORTS
# ============================================================================

import streamlit as st
import os
import io
import re
import json
import time
import base64
import tempfile
import shutil
import zipfile
import hashlib
import logging
from datetime import datetime, timedelta
from pathlib import Path
from typing import List, Dict, Tuple, Any, Optional, Union
import asyncio
from concurrent.futures import ThreadPoolExecutor
from functools import lru_cache, wraps

# Document processing
import pdfplumber
import pytesseract
from pdf2image import convert_from_path
from docx import Document as DocxDocument
from docx.shared import RGBColor as DocxRGBColor, Pt as DocxPt, Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd

# NLP and ML
from sklearn.feature_extraction.text import CountVectorizer
from sklearn.decomposition import LatentDirichletAllocation
import numpy as np

# OpenAI API
from openai import OpenAI
from tenacity import retry, wait_exponential, stop_after_attempt, retry_if_exception_type

# ============================================================================
# CONSTANTS & CONFIGURATION
# ============================================================================

# Version information
VERSION = "3.2"
BUILD_DATE = "2025-01-31"
ORGANIZATION = "QUINNS Training Services Limited"

# API Configuration
DEFAULT_MODEL = "gpt-4o"
MAX_TOKENS = 4000
TEMPERATURE = 0.7

# File constraints
MAX_FILE_SIZE_MB = 50
SUPPORTED_FORMATS = ['pdf', 'docx', 'pptx', 'txt', 'csv', 'xlsx', 'md']

# Training duration to slide mapping
DURATION_TO_SLIDES = {
    "30 minutes": 8,
    "1 hour": 12,
    "2 hours": 20,
    "half day": 30,
    "1 day": 50,
    "2 days": 80,
    "3 days": 120,
    "1 week": 200,
    "2 weeks": 360,
    "1 month": 720
}

# Module recommendations based on slides
MODULE_RECOMMENDATIONS = {
    (0, 50): 5,
    (51, 120): 6,
    (121, 200): 8,
    (201, 360): 10,
    (361, 720): 12,
    (721, 9999): 15
}

# Template options
TEMPLATE_OPTIONS = [
    "Corporate - Professional",
    "Technical - Detailed",
    "Compliance - Regulatory",
    "Sales - Persuasive",
    "Academic - Educational",
    "Workshop - Interactive"
]

# Phase identifiers
PHASE_1 = "PHASE_1"
PHASE_2 = "PHASE_2"
PHASE_3 = "PHASE_3"
PHASE_4 = "PHASE_4"

# Logging configuration
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('training_generator.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("EnhancedTrainingGenerator")

# Create output directories
OUTPUT_DIR = Path("output")
OUTPUT_DIR.mkdir(exist_ok=True)
CACHE_DIR = Path("cache")
CACHE_DIR.mkdir(exist_ok=True)
TEMP_DIR = Path("temp")
TEMP_DIR.mkdir(exist_ok=True)

# ============================================================================
# COMPLETE THEME SYSTEM - DARK AND LIGHT MODE WITH ALL FIXES
# ============================================================================

def get_theme_css(theme: str = "dark") -> str:
    """Generate CSS based on selected theme."""
    
    if theme == "light":
        return """
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700;800&display=swap');
    
    :root {
        --primary-color: #667eea;
        --secondary-color: #764ba2;
        --accent-pink: #f093fb;
        --accent-blue: #4facfe;
        --success-color: #10b981;
        --warning-color: #f59e0b;
        --error-color: #ef4444;
        --text-primary: #1e293b;
        --text-secondary: #334155;
        --text-muted: #64748b;
        --bg-primary: #ffffff;
        --bg-secondary: #f8fafc;
        --glass-bg: rgba(255, 255, 255, 0.9);
        --glass-border: rgba(0, 0, 0, 0.1);
        --shadow-color: rgba(0, 0, 0, 0.1);
    }
    
    * {
        font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
    }
    
    .main {
        background: transparent;
    }
    
    .stApp {
        background: linear-gradient(135deg, #f0f9ff 0%, #e0e7ff 25%, #ede9fe 50%, #fae8ff 75%, #fef3c7 100%);
        background-size: 200% 200%;
        animation: gradient 15s ease infinite;
        min-height: 100vh;
    }
    
    @keyframes gradient {
        0% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
        100% { background-position: 0% 50%; }
    }
    
    /* Text visibility */
    h1, h2, h3, h4, h5, h6 {
        color: var(--text-primary) !important;
        font-weight: 700 !important;
    }
    
    p, span, div:not(.stApp), label, li, td, th, a {
        color: var(--text-secondary) !important;
    }
    
    strong, b {
        color: var(--text-primary) !important;
        font-weight: 700 !important;
    }
    
    .stMarkdown, .stText {
        color: var(--text-primary) !important;
    }
    
    .caption, small, .stCaption {
        color: var(--text-muted) !important;
    }
    
    /* Input fields */
    .stTextInput input, 
    .stTextArea textarea, 
    .stNumberInput input,
    .stSelectbox select {
        color: var(--text-primary) !important;
        background-color: var(--bg-primary) !important;
        border: 2px solid var(--glass-border) !important;
        border-radius: 10px !important;
        padding: 12px 16px !important;
    }
    
    .stTextInput input:focus, 
    .stTextArea textarea:focus,
    .stNumberInput input:focus,
    .stSelectbox select:focus {
        border-color: var(--primary-color) !important;
        box-shadow: 0 0 0 3px rgba(102, 126, 234, 0.2) !important;
    }
    
    /* Number input controls - VISIBLE */
    .stNumberInput button {
        color: var(--primary-color) !important;
        background: var(--bg-secondary) !important;
        border: 1px solid var(--glass-border) !important;
        font-weight: 700 !important;
        font-size: 18px !important;
    }
    
    .stNumberInput button:hover {
        background: var(--primary-color) !important;
        color: white !important;
    }
    
    /* Input labels */
    .stTextInput label,
    .stTextArea label,
    .stNumberInput label,
    .stSelectbox label,
    .stSlider label,
    .stCheckbox label {
        color: var(--text-primary) !important;
        font-weight: 600 !important;
    }
    
    /* Selectbox/Dropdown - Light Theme */
    .stSelectbox > div > div {
        background-color: var(--bg-primary) !important;
        border: 2px solid var(--glass-border) !important;
    }
    
    .stSelectbox [data-baseweb="select"] > div,
    .stSelectbox [role="option"] {
        color: var(--text-primary) !important;
        background-color: var(--bg-primary) !important;
    }
    
    /* Info boxes */
    .stInfo {
        background: rgba(59, 130, 246, 0.1) !important;
        border-left: 4px solid #3b82f6 !important;
        color: var(--text-primary) !important;
    }
    
    .stSuccess {
        background: rgba(16, 185, 129, 0.1) !important;
        border-left: 4px solid var(--success-color) !important;
        color: var(--text-primary) !important;
    }
    
    .stWarning {
        background: rgba(245, 158, 11, 0.1) !important;
        border-left: 4px solid var(--warning-color) !important;
        color: var(--text-primary) !important;
    }
    
    .stError {
        background: rgba(239, 68, 68, 0.1) !important;
        border-left: 4px solid var(--error-color) !important;
        color: var(--text-primary) !important;
    }
    
    /* Expanders */
    .streamlit-expanderHeader {
        background: var(--bg-primary) !important;
        border: 2px solid var(--glass-border) !important;
        color: var(--text-primary) !important;
        font-weight: 600 !important;
    }
    
    .streamlit-expanderHeader:hover {
        background: var(--bg-secondary) !important;
        border-color: var(--primary-color) !important;
    }
    
    .streamlit-expanderContent {
        background: var(--bg-secondary) !important;
        border: 2px solid var(--glass-border) !important;
        border-top: none !important;
        color: var(--text-primary) !important;
    }
    
    /* Metrics */
    .stMetric {
        background: var(--bg-primary) !important;
        border: 2px solid var(--glass-border) !important;
        border-radius: 12px !important;
        padding: 20px !important;
    }
    
    .stMetric:hover {
        border-color: var(--primary-color) !important;
        box-shadow: 0 4px 12px var(--shadow-color) !important;
    }
    
    .stMetric label {
        color: var(--text-muted) !important;
    }
    
    .stMetric [data-testid="stMetricValue"] {
        color: var(--text-primary) !important;
        font-weight: 800 !important;
    }
    
    /* Buttons */
    .stButton button {
        background: linear-gradient(135deg, var(--primary-color), var(--secondary-color)) !important;
        color: white !important;
        border: none !important;
        border-radius: 12px !important;
        padding: 12px 24px !important;
        font-weight: 600 !important;
        box-shadow: 0 4px 12px rgba(102, 126, 234, 0.3) !important;
    }
    
    .stButton button:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 6px 20px rgba(102, 126, 234, 0.5) !important;
    }
    
    /* Download buttons - ALWAYS VISIBLE */
    .stDownloadButton button {
        background: linear-gradient(135deg, #10b981, #059669) !important;
        color: white !important;
        border: none !important;
        font-weight: 600 !important;
        box-shadow: 0 4px 12px rgba(16, 185, 129, 0.3) !important;
    }
    
    .stDownloadButton button:before {
        content: "📥 ";
        font-size: 20px;
        margin-right: 8px;
    }
    
    .stDownloadButton button:hover {
        background: linear-gradient(135deg, #059669, #047857) !important;
        transform: translateY(-2px) !important;
    }
    
    /* File uploader - DARK THEME WITH BLACK TEXT */
    .stFileUploader {
        background: rgba(248, 250, 252, 0.95) !important;  /* Light background */
        border: 2px dashed rgba(102, 126, 234, 0.5) !important;
        border-radius: 16px !important;
        padding: 32px !important;
    }
    
    .stFileUploader:hover {
        border-color: var(--primary-color) !important;
        background: rgba(255, 255, 255, 0.98) !important;
    }
    
    /* File uploader text - ALL BLACK */
    .stFileUploader label,
    .stFileUploader span,
    .stFileUploader p,
    .stFileUploader div {
        color: #0f172a !important;  /* Very dark/black */
        font-weight: 500 !important;
    }
    
    /* File uploader label (main text) */
    .stFileUploader > label > div {
        color: #020617 !important;  /* Pure black */
        font-weight: 600 !important;
    }
    
    /* Drag and drop text */
    .stFileUploader [data-testid="stFileUploaderDropzone"] {
        background: white !important;
    }
    
    .stFileUploader [data-testid="stFileUploaderDropzone"] span,
    .stFileUploader [data-testid="stFileUploaderDropzone"] p,
    .stFileUploader [data-testid="stFileUploaderDropzone"] small {
        color: #1e293b !important;  /* Dark slate */
        font-weight: 500 !important;
    }
    
    /* File uploader instructions text */
    .stFileUploader small {
        color: #334155 !important;  /* Medium dark */
    }
    
    /* Progress bar */
    .stProgress > div > div {
        background: linear-gradient(90deg, var(--primary-color), var(--accent-pink)) !important;
    }
    
    .stProgress > div {
        background: rgba(0, 0, 0, 0.1) !important;
    }
    
    /* Tabs */
    .stTabs [data-baseweb="tab-list"] {
        background: var(--bg-primary);
        border: 2px solid var(--glass-border);
        border-radius: 12px;
        padding: 8px;
    }
    
    .stTabs [data-baseweb="tab"] {
        color: var(--text-muted) !important;
        font-weight: 600 !important;
    }
    
    .stTabs [aria-selected="true"] {
        background: linear-gradient(135deg, var(--primary-color), var(--secondary-color)) !important;
        color: white !important;
    }
    
    /* Sidebar */
    section[data-testid="stSidebar"] {
        background: var(--bg-secondary) !important;
        border-right: 2px solid var(--glass-border) !important;
    }
    
    section[data-testid="stSidebar"] * {
        color: var(--text-secondary) !important;
    }
    
    section[data-testid="stSidebar"] h1,
    section[data-testid="stSidebar"] h2,
    section[data-testid="stSidebar"] h3 {
        color: var(--text-primary) !important;
    }
    
    /* Phase badges */
    .phase-badge {
        display: inline-block;
        padding: 10px 20px;
        border-radius: 25px;
        font-weight: 600;
        font-size: 14px;
        margin: 8px 0;
        text-align: center;
    }
    
    .phase-active {
        background: linear-gradient(135deg, var(--primary-color), var(--secondary-color));
        color: white !important;
        box-shadow: 0 4px 15px rgba(102, 126, 234, 0.4);
    }
    
    .phase-complete {
        background: linear-gradient(135deg, var(--success-color), #059669);
        color: white !important;
    }
    
    .phase-pending {
        background: var(--bg-secondary);
        color: var(--text-muted) !important;
        border: 2px solid var(--glass-border);
    }
    
    /* Importance badges - COLOR CODED */
    .importance-high {
        background: linear-gradient(135deg, #ef4444, #dc2626) !important;
        color: white !important;
        padding: 4px 12px;
        border-radius: 20px;
        font-weight: 600;
        display: inline-block;
    }
    
    .importance-medium {
        background: linear-gradient(135deg, #f59e0b, #d97706) !important;
        color: white !important;
        padding: 4px 12px;
        border-radius: 20px;
        font-weight: 600;
        display: inline-block;
    }
    
    .importance-low {
        background: linear-gradient(135deg, #3b82f6, #2563eb) !important;
        color: white !important;
        padding: 4px 12px;
        border-radius: 20px;
        font-weight: 600;
        display: inline-block;
    }
    
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    
    .block-container {
        padding-top: 2rem !important;
        max-width: 1200px !important;
    }
</style>
"""
    else:  # Dark theme
        return """
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700;800&display=swap');
    
    :root {
        --primary-color: #667eea;
        --secondary-color: #764ba2;
        --accent-pink: #f093fb;
        --accent-blue: #4facfe;
        --success-color: #10b981;
        --warning-color: #fbbf24;
        --error-color: #ef4444;
        --text-primary: #ffffff;
        --text-secondary: rgba(255, 255, 255, 0.95);
        --text-muted: rgba(255, 255, 255, 0.7);
        --glass-bg: rgba(30, 27, 75, 0.6);
        --glass-border: rgba(255, 255, 255, 0.15);
        --shadow-color: rgba(102, 126, 234, 0.3);
    }
    
    * {
        font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
    }
    
    .main {
        background: transparent;
    }
    
    .stApp {
        background: linear-gradient(135deg, #0f172a 0%, #1e1b4b 25%, #581c87 50%, #1e3a8a 75%, #0f172a 100%);
        background-size: 200% 200%;
        animation: gradient 10s ease infinite;
        min-height: 100vh;
    }
    
    @keyframes gradient {
        0% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
        100% { background-position: 0% 50%; }
    }
    
    /* Text visibility */
    h1, h2, h3, h4, h5, h6 {
        color: var(--text-primary) !important;
        font-weight: 700 !important;
        text-shadow: 0 2px 8px rgba(0, 0, 0, 0.5);
    }
    
    p, span, div:not(.stApp), label, li, td, th, a {
        color: var(--text-secondary) !important;
    }
    
    strong, b {
        color: var(--text-primary) !important;
        font-weight: 700 !important;
    }
    
    .stMarkdown, .stText {
        color: var(--text-primary) !important;
    }
    
    .caption, small, .stCaption {
        color: var(--text-muted) !important;
    }
    
    /* Input fields - WHITE BACKGROUND WITH BLACK TEXT */
    .stTextInput input, 
    .stTextArea textarea, 
    .stNumberInput input {
        color: #0f172a !important;
        background-color: white !important;
        border: 2px solid rgba(102, 126, 234, 0.5) !important;
        border-radius: 10px !important;
        padding: 12px 16px !important;
        font-weight: 500 !important;
    }
    
    .stTextInput input:focus, 
    .stTextArea textarea:focus,
    .stNumberInput input:focus {
        background-color: white !important;
        border-color: var(--primary-color) !important;
        box-shadow: 0 0 0 3px rgba(102, 126, 234, 0.2) !important;
        color: #0f172a !important;
    }
    
    .stTextInput input::placeholder,
    .stTextArea textarea::placeholder {
        color: #94a3b8 !important;
    }
    
    /* Number input controls - HIGHLY VISIBLE */
    .stNumberInput button {
        color: #667eea !important;
        background: rgba(102, 126, 234, 0.2) !important;
        border: 2px solid #667eea !important;
        font-weight: 900 !important;
        font-size: 20px !important;
        border-radius: 6px !important;
    }
    
    .stNumberInput button:hover {
        background: linear-gradient(135deg, #667eea, #764ba2) !important;
        color: white !important;
        transform: scale(1.1);
    }
    
    /* Input labels - WHITE */
    .stTextInput label,
    .stTextArea label,
    .stNumberInput label,
    .stSlider label,
    .stCheckbox label {
        color: var(--text-primary) !important;
        font-weight: 600 !important;
    }
    
    /* ========== SELECTBOX/DROPDOWN - WHITE BACKGROUND WITH BLACK TEXT ========== */
    
    .stSelectbox {
        color: #0f172a !important;
    }
    
    .stSelectbox label {
        color: #ffffff !important;
        font-weight: 600 !important;
    }
    
    /* Selectbox container - WHITE BACKGROUND */
    .stSelectbox > div > div {
        background-color: white !important;
        border: 2px solid rgba(102, 126, 234, 0.5) !important;
        border-radius: 10px !important;
    }
    
    /* Selected value display - BLACK TEXT */
    .stSelectbox [data-baseweb="select"] {
        background-color: white !important;
    }
    
    .stSelectbox [data-baseweb="select"] > div,
    .stSelectbox [data-baseweb="select"] span,
    .stSelectbox [data-baseweb="select"] input {
        color: #0f172a !important;
        background-color: white !important;
        font-weight: 600 !important;
    }
    
    /* Dropdown arrow - BLACK */
    .stSelectbox svg {
        fill: #0f172a !important;
    }
    
    /* Dropdown menu popover - WHITE BACKGROUND */
    .stSelectbox [data-baseweb="popover"],
    .stSelectbox [data-baseweb="menu"],
    .stSelectbox [role="listbox"],
    div[role="listbox"],
    ul[role="listbox"] {
        background-color: white !important;
        border: 2px solid rgba(102, 126, 234, 0.5) !important;
        border-radius: 10px !important;
        box-shadow: 0 8px 20px rgba(0, 0, 0, 0.3) !important;
    }
    
    /* Dropdown options - BLACK TEXT */
    .stSelectbox [role="option"],
    .stSelectbox li,
    div[role="listbox"] li,
    ul[role="listbox"] li {
        color: #0f172a !important;
        background-color: white !important;
        font-weight: 500 !important;
        padding: 12px 16px !important;
    }
    
    /* Dropdown option hover - LIGHT PURPLE */
    .stSelectbox [role="option"]:hover,
    .stSelectbox li:hover,
    div[role="listbox"] li:hover,
    ul[role="listbox"] li:hover {
        background-color: rgba(102, 126, 234, 0.15) !important;
        color: #0f172a !important;
        font-weight: 600 !important;
    }
    
    /* Selected option - PURPLE BACKGROUND */
    .stSelectbox [aria-selected="true"],
    div[role="listbox"] li[aria-selected="true"],
    ul[role="listbox"] li[aria-selected="true"] {
        background-color: rgba(102, 126, 234, 0.3) !important;
        color: #0f172a !important;
        font-weight: 700 !important;
    }
    
    /* ========== END SELECTBOX STYLING ========== */
    
    /* Info boxes - LIGHT BACKGROUND WITH DARK TEXT */
    .stInfo {
        background: linear-gradient(135deg, rgba(239, 246, 255, 0.98), rgba(224, 242, 254, 0.98)) !important;
        border-left: 4px solid #3b82f6 !important;
        border-radius: 8px !important;
        padding: 16px !important;
    }
    
    .stInfo, .stInfo * {
        color: #1e3a8a !important;
        font-weight: 500 !important;
    }
    
    .stSuccess {
        background: linear-gradient(135deg, rgba(236, 253, 245, 0.98), rgba(209, 250, 229, 0.98)) !important;
        border-left: 4px solid var(--success-color) !important;
        border-radius: 8px !important;
        padding: 16px !important;
    }
    
    .stSuccess, .stSuccess * {
        color: #065f46 !important;
        font-weight: 500 !important;
    }
    
    .stWarning {
        background: linear-gradient(135deg, rgba(254, 252, 232, 0.98), rgba(254, 249, 195, 0.98)) !important;
        border-left: 4px solid var(--warning-color) !important;
        border-radius: 8px !important;
        padding: 16px !important;
    }
    
    .stWarning, .stWarning * {
        color: #92400e !important;
        font-weight: 500 !important;
    }
    
    .stError {
        background: linear-gradient(135deg, rgba(254, 242, 242, 0.98), rgba(254, 226, 226, 0.98)) !important;
        border-left: 4px solid var(--error-color) !important;
        border-radius: 8px !important;
        padding: 16px !important;
    }
    
    .stError, .stError * {
        color: #7f1d1d !important;
        font-weight: 500 !important;
    }
    
    /* Expanders */
    .streamlit-expanderHeader {
        background: rgba(30, 27, 75, 0.6) !important;
        border: 1px solid var(--glass-border) !important;
        color: white !important;
        font-weight: 600 !important;
    }
    
    .streamlit-expanderHeader:hover {
        background: rgba(102, 126, 234, 0.2) !important;
        border-color: var(--primary-color) !important;
    }
    
    .streamlit-expanderContent {
        background: rgba(30, 27, 75, 0.4) !important;
        border: 1px solid var(--glass-border) !important;
        border-top: none !important;
        color: white !important;
    }
    
    /* Metrics */
    .stMetric {
        background: rgba(30, 27, 75, 0.6) !important;
        border: 1px solid var(--glass-border) !important;
        border-radius: 12px !important;
        padding: 20px !important;
    }
    
    .stMetric:hover {
        border-color: var(--primary-color) !important;
        box-shadow: 0 8px 25px var(--shadow-color) !important;
    }
    
    .stMetric label {
        color: var(--text-muted) !important;
    }
    
    .stMetric [data-testid="stMetricValue"] {
        color: var(--text-primary) !important;
        font-weight: 800 !important;
    }
    
    /* Buttons */
    .stButton button {
        background: linear-gradient(135deg, var(--primary-color), var(--secondary-color)) !important;
        color: white !important;
        border: none !important;
        border-radius: 12px !important;
        padding: 12px 24px !important;
        font-weight: 600 !important;
        box-shadow: 0 4px 15px rgba(102, 126, 234, 0.4) !important;
    }
    
    .stButton button:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 6px 20px rgba(102, 126, 234, 0.6) !important;
    }
    
    /* Download buttons - ALWAYS VISIBLE */
    .stDownloadButton button {
        background: linear-gradient(135deg, #10b981, #059669) !important;
        color: white !important;
        border: none !important;
        font-weight: 600 !important;
        box-shadow: 0 4px 15px rgba(16, 185, 129, 0.4) !important;
    }
    
    .stDownloadButton button:before {
        content: "📥 ";
        font-size: 20px;
        margin-right: 8px;
    }
    
    .stDownloadButton button:hover {
        background: linear-gradient(135deg, #059669, #047857) !important;
        transform: translateY(-2px) !important;
        box-shadow: 0 6px 25px rgba(16, 185, 129, 0.6) !important;
    }
    
    /* File uploader - LIGHT BACKGROUND WITH BLACK TEXT */
    .stFileUploader {
        background: linear-gradient(135deg, rgba(248, 250, 252, 0.98), rgba(241, 245, 249, 0.98)) !important;
        border: 2px dashed rgba(102, 126, 234, 0.6) !important;
        border-radius: 16px !important;
        padding: 32px !important;
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.1) !important;
    }
    
    .stFileUploader:hover {
        border-color: #667eea !important;
        background: rgba(255, 255, 255, 0.99) !important;
        box-shadow: 0 6px 20px rgba(102, 126, 234, 0.3) !important;
    }
    
    /* File uploader label - WHITE TEXT */
    .stFileUploader > label {
        color: #ffffff !important;
        font-weight: 600 !important;
        font-size: 1rem !important;
    }
    
    /* File uploader dropzone - WHITE BACKGROUND */
    .stFileUploader [data-testid="stFileUploaderDropzone"] {
        background: white !important;
        border-radius: 12px !important;
    }
    
    /* All text inside file uploader dropzone - BLACK */
    .stFileUploader [data-testid="stFileUploaderDropzone"] *,
    .stFileUploader [data-testid="stFileUploaderDropzone"] span,
    .stFileUploader [data-testid="stFileUploaderDropzone"] p,
    .stFileUploader [data-testid="stFileUploaderDropzone"] div,
    .stFileUploader [data-testid="stFileUploaderDropzone"] small {
        color: #1e293b !important;
        font-weight: 500 !important;
    }
    
    /* File name when uploaded */
    .stFileUploader [data-testid="stFileUploaderFileName"] {
        color: #0f172a !important;
        font-weight: 600 !important;
    }
    
    /* Browse files button */
    .stFileUploader button {
        color: white !important;
        background: linear-gradient(135deg, #667eea, #764ba2) !important;
    }
    
    /* Progress bar */
    .stProgress > div > div {
        background: linear-gradient(90deg, var(--primary-color), var(--accent-pink)) !important;
    }
    
    .stProgress > div {
        background: rgba(255, 255, 255, 0.1) !important;
    }
    
    /* Tabs */
    .stTabs [data-baseweb="tab-list"] {
        background: rgba(30, 27, 75, 0.6);
        border: 1px solid var(--glass-border);
        border-radius: 12px;
        padding: 8px;
    }
    
    .stTabs [data-baseweb="tab"] {
        color: var(--text-muted) !important;
        font-weight: 600 !important;
    }
    
    .stTabs [aria-selected="true"] {
        background: linear-gradient(135deg, var(--primary-color), var(--secondary-color)) !important;
        color: white !important;
    }
    
    /* Sidebar - LIGHT BACKGROUND WITH BLACK TEXT */
    section[data-testid="stSidebar"] {
        background: rgba(248, 250, 252, 0.98) !important;
        border-right: 1px solid var(--glass-border) !important;
    }
    
    section[data-testid="stSidebar"] * {
        color: #0f172a !important;
    }
    
    section[data-testid="stSidebar"] h1,
    section[data-testid="stSidebar"] h2,
    section[data-testid="stSidebar"] h3 {
        color: #020617 !important;
        font-weight: 700 !important;
    }
    
    section[data-testid="stSidebar"] label,
    section[data-testid="stSidebar"] p,
    section[data-testid="stSidebar"] span,
    section[data-testid="stSidebar"] .stCaption,
    section[data-testid="stSidebar"] div {
        color: #1e293b !important;
    }
    
    /* Sidebar input fields */
    section[data-testid="stSidebar"] input,
    section[data-testid="stSidebar"] select,
    section[data-testid="stSidebar"] textarea {
        color: #0f172a !important;
        background-color: white !important;
    }
    
    /* Keep sidebar button text white */
    section[data-testid="stSidebar"] .stButton button {
        color: white !important;
    }
    
    /* Phase badges */
    .phase-badge {
        display: inline-block;
        padding: 10px 20px;
        border-radius: 25px;
        font-weight: 600;
        font-size: 14px;
        margin: 8px 0;
        text-align: center;
    }
    
    .phase-active {
        background: linear-gradient(135deg, var(--primary-color), var(--secondary-color));
        color: white !important;
        box-shadow: 0 4px 15px rgba(102, 126, 234, 0.5);
    }
    
    .phase-complete {
        background: linear-gradient(135deg, var(--success-color), #059669);
        color: white !important;
    }
    
    .phase-pending {
        background: rgba(30, 27, 75, 0.6);
        color: var(--text-muted) !important;
        border: 1px solid var(--glass-border);
    }
    
    /* Importance badges - COLOR CODED */
    .importance-high {
        background: linear-gradient(135deg, #ef4444, #dc2626) !important;
        color: white !important;
        padding: 6px 16px;
        border-radius: 20px;
        font-weight: 700;
        display: inline-block;
        box-shadow: 0 2px 8px rgba(239, 68, 68, 0.4);
    }
    
    .importance-medium {
        background: linear-gradient(135deg, #f59e0b, #d97706) !important;
        color: white !important;
        padding: 6px 16px;
        border-radius: 20px;
        font-weight: 700;
        display: inline-block;
        box-shadow: 0 2px 8px rgba(245, 158, 11, 0.4);
    }
    
    .importance-low {
        background: linear-gradient(135deg, #3b82f6, #2563eb) !important;
        color: white !important;
        padding: 6px 16px;
        border-radius: 20px;
        font-weight: 700;
        display: inline-block;
        box-shadow: 0 2px 8px rgba(59, 130, 246, 0.4);
    }
    
    /* FORCE DROPDOWN TEXT BLACK - GLOBAL OVERRIDE */
    [data-baseweb="select"] *:not(svg),
    [data-baseweb="popover"] *:not(svg),
    [data-baseweb="menu"] *:not(svg),
    [role="listbox"] *:not(svg),
    div[role="listbox"] *:not(svg),
    ul[role="listbox"] *:not(svg),
    li[role="option"] *:not(svg) {
        color: #0f172a !important;
    }
    
    [data-baseweb="select"],
    [data-baseweb="popover"],
    [data-baseweb="menu"],
    [role="listbox"],
    div[role="listbox"],
    ul[role="listbox"] {
        background-color: white !important;
    }
    
    /* Keep selectbox label white */
    .stSelectbox > label,
    .stSelectbox > label * {
        color: #ffffff !important;
    }
    
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    
    .block-container {
        padding-top: 2rem !important;
        max-width: 1200px !important;
    }
</style>
"""

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def get_file_hash(content: bytes) -> str:
    """Generate hash for file caching."""
    return hashlib.md5(content).hexdigest()

def validate_api_key(api_key: str) -> bool:
    """Validate OpenAI API key format."""
    if not api_key:
        return False
    return api_key.startswith('sk-') and len(api_key) > 20

def calculate_reading_time(word_count: int) -> int:
    """Calculate estimated reading time in minutes (250 words/min)."""
    return max(1, word_count // 250)

def get_recommended_modules(slide_count: int) -> int:
    """Get recommended module count based on slide count."""
    for (min_slides, max_slides), modules in MODULE_RECOMMENDATIONS.items():
        if min_slides <= slide_count <= max_slides:
            return modules
    return 8

def format_duration(minutes: int) -> str:
    """Format duration in minutes to human-readable string."""
    hours = minutes // 60
    mins = minutes % 60
    
    if hours == 0:
        return f"{mins}m"
    elif mins == 0:
        return f"{hours}h"
    else:
        return f"{hours}h {mins}m"

def clean_json_response(response_text: str) -> str:
    """Clean JSON response from API."""
    text = response_text.strip()
    
    if text.startswith("```json"):
        text = text[7:]
    elif text.startswith("```"):
        text = text[3:]
    
    if text.endswith("```"):
        text = text[:-3]
    
    return text.strip()

def safe_json_load(text: str) -> Dict:
    """Safely load JSON from API response."""
    text = text.strip()
    text = text.removeprefix("```json").removeprefix("```").strip()
    text = text.removesuffix("```").strip()
    
    if not text.startswith('{') and not text.startswith('['):
        import re
        json_match = re.search(r'(\{.*\}|\[.*\])', text, re.DOTALL)
        if json_match:
            text = json_match.group(1)
    
    return json.loads(text)

def chunk_text(text: str, max_chars: int = 80000) -> List[str]:
    """Split text into chunks safe for GPT context window."""
    if len(text) <= max_chars:
        return [text]
    
    chunks = []
    for i in range(0, len(text), max_chars):
        chunks.append(text[i:i + max_chars])
    
    logger.info(f"Chunked {len(text)} chars into {len(chunks)} chunks")
    return chunks

def validate_file_upload(uploaded_file) -> Tuple[bool, str]:
    """Validate uploaded file."""
    if uploaded_file is None:
        return False, "No file uploaded"
    
    file_size_mb = uploaded_file.size / (1024 * 1024)
    if file_size_mb > MAX_FILE_SIZE_MB:
        return False, f"File size ({file_size_mb:.1f}MB) exceeds maximum ({MAX_FILE_SIZE_MB}MB)"
    
    file_extension = uploaded_file.name.split('.')[-1].lower()
    if file_extension not in SUPPORTED_FORMATS:
        return False, f"Unsupported format: {file_extension}"
    
    return True, "Valid file"

def load_api_key_from_env() -> str:
    """Load API key from Streamlit secrets.toml file."""
    try:
        api_key = st.secrets["OPENAI_API_KEY"]
        if api_key:
            logger.info("API key loaded from secrets.toml")
        return api_key
    except (KeyError, FileNotFoundError):
        logger.warning("API key not found in secrets.toml")
        return ""

def initialize_session_state():
    """Initialize all session state variables."""
    defaults = {
        'api_key': load_api_key_from_env(),  # Auto-load from secrets.toml
        'theme': 'light',  # Default theme
        'current_phase': PHASE_1,
        'phase_completed': {
            PHASE_1: False,
            PHASE_2: False,
            PHASE_3: False,
            PHASE_4: False
        },
        'uploaded_file': None,
        'processed_content': None,
        'full_source_text': None,
        'extracted_topics': None,
        'generated_outline': None,
        'edited_topics': None,
        'edited_outline': None,
        'generated_slides': None,
        'final_documents': {},
        'target_modules': 8,
        'target_slides': 50,
        'training_duration': '1 day',
        'extended_mode': False,
        'custom_slide_count': None,
        'template': 'Corporate - Professional',
        'include_assessments': True,
        'include_activities': True,
        'detailed_analytics': False,
        'enhanced_trainer_guide': True,
        'generation_progress': {},
        'error_log': []
    }
    
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

def render_importance_badge(importance: str) -> str:
    """Render color-coded importance badge."""
    importance_lower = importance.lower()
    if importance_lower == 'high':
        return f'<span class="importance-high">🔴 HIGH</span>'
    elif importance_lower == 'medium':
        return f'<span class="importance-medium">🟡 MEDIUM</span>'
    elif importance_lower == 'low':
        return f'<span class="importance-low">🔵 LOW</span>'
    else:
        return f'<span class="importance-medium">🟡 MEDIUM</span>'

# ============================================================================
# DOCUMENT PROCESSOR CLASS (keeping same as before for brevity)
# ============================================================================

class DocumentProcessor:
    """Handles extraction and processing of various document formats."""
    
    def __init__(self):
        self.supported_formats = SUPPORTED_FORMATS
        logger.info("DocumentProcessor initialized")
    
    def process_file(self, uploaded_file) -> Dict[str, Any]:
        """Process uploaded file and extract content."""
        try:
            file_extension = uploaded_file.name.split('.')[-1].lower()
            logger.info(f"Processing {file_extension} file: {uploaded_file.name}")
            
            temp_path = TEMP_DIR / uploaded_file.name
            with open(temp_path, 'wb') as f:
                f.write(uploaded_file.getbuffer())
            
            if file_extension == 'pdf':
                result = self._process_pdf(temp_path)
            elif file_extension == 'docx':
                result = self._process_docx(temp_path)
            elif file_extension == 'pptx':
                result = self._process_pptx(temp_path)
            elif file_extension == 'txt':
                result = self._process_txt(temp_path)
            elif file_extension == 'md':
                result = self._process_markdown(temp_path)
            elif file_extension in ['csv', 'xlsx']:
                result = self._process_spreadsheet(temp_path, file_extension)
            else:
                raise ValueError(f"Unsupported format: {file_extension}")
            
            result['format'] = file_extension
            result['filename'] = uploaded_file.name
            result['processed_at'] = datetime.now().isoformat()
            result['word_count'] = len(result['text'].split())
            result['bullets'] = self._extract_bullets(result['text'])
            
            if temp_path.exists():
                temp_path.unlink()
            
            logger.info(f"Successfully processed {uploaded_file.name}: {result['word_count']} words")
            return result
            
        except Exception as e:
            logger.error(f"Error processing file: {str(e)}")
            raise
    
    def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PDF file."""
        text_content = []
        pages_data = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text() or ""
                    
                    if not page_text.strip():
                        try:
                            images = convert_from_path(file_path, first_page=i, last_page=i)
                            if images:
                                page_text = pytesseract.image_to_string(images[0])
                        except Exception as ocr_error:
                            logger.warning(f"OCR failed for page {i}: {str(ocr_error)}")
                    
                    text_content.append(page_text)
                    pages_data.append({
                        'page_number': i,
                        'text': page_text,
                        'bullets': self._extract_bullets(page_text)
                    })
        
        except Exception as e:
            logger.error(f"PDF processing error: {str(e)}")
            raise
        
        return {
            'text': '\n\n'.join(text_content),
            'pages': pages_data,
            'page_count': len(pages_data)
        }
    
    def _process_docx(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from Word document."""
        doc = DocxDocument(file_path)
        
        paragraphs = []
        pages_data = []
        current_page_text = []
        page_num = 1
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
                current_page_text.append(text)
                
                if len(' '.join(current_page_text).split()) > 500:
                    page_text = '\n'.join(current_page_text)
                    pages_data.append({
                        'page_number': page_num,
                        'text': page_text,
                        'bullets': self._extract_bullets(page_text)
                    })
                    current_page_text = []
                    page_num += 1
        
        if current_page_text:
            page_text = '\n'.join(current_page_text)
            pages_data.append({
                'page_number': page_num,
                'text': page_text,
                'bullets': self._extract_bullets(page_text)
            })
        
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join([cell.text for cell in row.cells])
                tables_text.append(row_text)
        
        full_text = '\n\n'.join(paragraphs)
        if tables_text:
            full_text += '\n\nTABLES:\n' + '\n'.join(tables_text)
        
        return {
            'text': full_text,
            'pages': pages_data,
            'page_count': len(pages_data)
        }
    
    def _process_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PowerPoint presentation."""
        prs = PptxPresentation(file_path)
        
        slides_data = []
        text_content = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text_parts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text)
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = ' | '.join([cell.text for cell in row.cells])
                        slide_text_parts.append(row_text)
            
            slide_text = '\n'.join(slide_text_parts)
            text_content.append(slide_text)
            
            slides_data.append({
                'slide_number': i,
                'text': slide_text,
                'bullets': self._extract_bullets(slide_text)
            })
        
        return {
            'text': '\n\n'.join(text_content),
            'pages': slides_data,
            'page_count': len(slides_data)
        }
    
    def _process_txt(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from text file."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        
        words = text.split()
        pages_data = []
        page_size = 1000
        
        for i in range(0, len(words), page_size):
            page_text = ' '.join(words[i:i+page_size])
            pages_data.append({
                'page_number': len(pages_data) + 1,
                'text': page_text,
                'bullets': self._extract_bullets(page_text)
            })
        
        return {
            'text': text,
            'pages': pages_data,
            'page_count': len(pages_data)
        }
    
    def _process_markdown(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from markdown file."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        
        sections = re.split(r'\n#+\s+', text)
        pages_data = []
        
        for i, section in enumerate(sections, 1):
            if section.strip():
                pages_data.append({
                    'page_number': i,
                    'text': section,
                    'bullets': self._extract_bullets(section)
                })
        
        return {
            'text': text,
            'pages': pages_data,
            'page_count': len(pages_data)
        }
    
    def _process_spreadsheet(self, file_path: Path, format: str) -> Dict[str, Any]:
        """Extract content from CSV or Excel file."""
        if format == 'csv':
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
        
        text_parts = []
        text_parts.append("HEADERS: " + ', '.join(df.columns.tolist()))
        
        for idx, row in df.iterrows():
            row_text = ' | '.join([f"{col}: {val}" for col, val in row.items()])
            text_parts.append(row_text)
        
        full_text = '\n'.join(text_parts)
        
        summary = f"\nSUMMARY:\nRows: {len(df)}\nColumns: {len(df.columns)}\n"
        summary += f"Columns: {', '.join(df.columns.tolist())}"
        
        return {
            'text': full_text + summary,
            'pages': [{'page_number': 1, 'text': full_text, 'bullets': []}],
            'page_count': 1,
            'dataframe': df.to_dict('records')
        }
    
    def _extract_bullets(self, text: str) -> List[str]:
        """Extract bullet points and key sentences from text."""
        bullets = []
        
        bullet_patterns = [
            r'^\s*[\-\*\•]\s+(.+)$',
            r'^\s*\d+[\.\)]\s+(.+)$',
            r'^\s*[a-zA-Z][\.\)]\s+(.+)$'
        ]
        
        for line in text.split('\n'):
            for pattern in bullet_patterns:
                match = re.match(pattern, line)
                if match:
                    bullets.append(match.group(1).strip())
                    break
        
        if not bullets:
            sentences = re.split(r'[.!?]+', text)
            bullets = [s.strip() for s in sentences if len(s.split()) >= 5][:10]
        
        return bullets[:20]

# ============================================================================
# AI CLASSES (TopicAnalyzer, OutlineGenerator, SlideGenerator, DocumentBuilder)
# Keeping these the same as before for brevity
# ============================================================================

class TopicAnalyzer:
    """Analyzes document content and extracts main topics using AI."""
    
    def __init__(self, api_key: str):
        self.client = OpenAI(api_key=api_key)
        self.model = DEFAULT_MODEL
        logger.info("TopicAnalyzer initialized")
    
    @retry(wait=wait_exponential(min=4, max=60), stop=stop_after_attempt(3))
    def extract_topics(self, content: str, num_topics: int = 8) -> List[Dict[str, Any]]:
        """Extract main topics from document content using AI."""
        try:
            logger.info(f"Extracting {num_topics} topics from content")
            
            chunks = chunk_text(content, 100000)
            summary = "\n\n".join([chunk[:5000] for chunk in chunks[:3]])
            
            logger.info(f"Using {len(summary)} chars of content for topic extraction")
            
            prompt = f"""Analyze this training document and extract {num_topics} distinct main topics.

DOCUMENT:
{summary}

Return JSON with this EXACT structure:
{{
  "topics": [
    {{
      "id": 1,
      "title": "Specific Topic Name",
      "description": "What this topic covers (2-3 sentences)",
      "key_concepts": ["concept1", "concept2", "concept3", "concept4", "concept5"],
      "importance": "high",
      "estimated_duration_minutes": 45
    }}
  ]
}}

Requirements:
- Extract {num_topics} topics
- Each topic must have 4-6 key concepts
- Importance: high/medium/low
- Duration: 20-90 minutes based on topic complexity
- Topics should be distinct and cover different aspects
- Return ONLY valid JSON, no additional text"""

            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert instructional designer analyzing training content. Return only valid JSON."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                response_format={"type": "json_object"},
                temperature=TEMPERATURE,
                max_tokens=MAX_TOKENS
            )
            
            response_text = response.choices[0].message.content
            result = safe_json_load(response_text)
            topics = result.get('topics', [])
            
            for i, topic in enumerate(topics, 1):
                if 'id' not in topic:
                    topic['id'] = i
                if 'importance' not in topic or topic['importance'] not in ['high', 'medium', 'low']:
                    topic['importance'] = 'medium'
                if 'estimated_duration_minutes' not in topic:
                    topic['estimated_duration_minutes'] = 45
            
            logger.info(f"Successfully extracted {len(topics)} topics")
            return topics
            
        except json.JSONDecodeError as e:
            logger.error(f"JSON parsing error: {str(e)}")
            logger.error(f"Response text: {response_text}")
            raise ValueError(f"Failed to parse API response as JSON: {str(e)}")
        
        except Exception as e:
            logger.error(f"Topic extraction error: {str(e)}")
            raise

class OutlineGenerator:
    """Generates structured training outline from topics."""
    
    def __init__(self, api_key: str):
        self.client = OpenAI(api_key=api_key)
        self.model = DEFAULT_MODEL
        logger.info("OutlineGenerator initialized")
    
    @retry(wait=wait_exponential(multiplier=1, min=4, max=60), stop=stop_after_attempt(5))
    def generate_outline(
        self,
        topics: List[Dict],
        target_modules: int,
        target_slides: int,
        duration: str,
        template: str = "Corporate - Professional"
    ) -> Dict[str, Any]:
        """Generate structured training outline."""
        try:
            logger.info(f"Generating outline: {target_modules} modules, {target_slides} slides")
            
            topics_summary = "\n".join([
                f"{i}. {t['title']}\n   Concepts: {', '.join(t['key_concepts'][:4])}\n   Duration: {t['estimated_duration_minutes']}min"
                for i, t in enumerate(topics, 1)
            ])
            
            slides_per_module = target_slides // target_modules
            
            prompt = f"""Create a detailed training outline with {target_modules} modules and approximately {target_slides} total slides.

TOPICS IDENTIFIED:
{topics_summary}

REQUIREMENTS:
- Duration: {duration}
- Modules: {target_modules}
- Total Slides: ~{target_slides} (approximately {slides_per_module} per module)
- Template Style: {template}

Return JSON with this EXACT structure:
{{
  "title": "Comprehensive Training Program Title",
  "description": "2-3 sentence program overview",
  "duration": "{duration}",
  "total_modules": {target_modules},
  "estimated_slides": {target_slides},
  "objectives": [
    "Primary learning objective 1",
    "Primary learning objective 2",
    "Primary learning objective 3",
    "Primary learning objective 4"
  ],
  "modules": [
    {{
      "id": 1,
      "title": "Module Title",
      "duration": "2 hours",
      "objectives": [
        "Module objective 1",
        "Module objective 2",
        "Module objective 3"
      ],
      "topics": ["topic1", "topic2", "topic3"],
      "key_points": [
        "Key point 1",
        "Key point 2",
        "Key point 3",
        "Key point 4",
        "Key point 5"
      ],
      "estimated_slides": {slides_per_module}
    }}
  ]
}}

Requirements:
- Create exactly {target_modules} modules
- Distribute {target_slides} slides evenly across modules
- Each module should have 3-4 objectives
- Each module should have 5-8 key points
- Modules should build progressively
- Return ONLY valid JSON"""

            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert instructional designer creating comprehensive training outlines. Return only valid JSON."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                response_format={"type": "json_object"},
                temperature=TEMPERATURE,
                max_tokens=MAX_TOKENS
            )
            
            response_text = response.choices[0].message.content
            outline = safe_json_load(response_text)
            
            for i, module in enumerate(outline.get('modules', []), 1):
                if 'id' not in module:
                    module['id'] = i
                if 'estimated_slides' not in module:
                    module['estimated_slides'] = slides_per_module
            
            actual_total = sum(m['estimated_slides'] for m in outline['modules'])
            if actual_total != target_slides:
                diff = target_slides - actual_total
                outline['modules'][0]['estimated_slides'] += diff
            
            outline['estimated_slides'] = target_slides
            outline['total_modules'] = len(outline['modules'])
            
            logger.info(f"Successfully generated outline with {len(outline['modules'])} modules")
            return outline
            
        except json.JSONDecodeError as e:
            logger.error(f"JSON parsing error: {str(e)}")
            logger.error(f"Response text: {response_text}")
            raise ValueError(f"Failed to parse API response as JSON: {str(e)}")
        
        except Exception as e:
            logger.error(f"Outline generation error: {str(e)}")
            raise

class SlideGenerator:
    """Generates individual slides for each module using AI."""
    
    def __init__(self, api_key: str):
        self.client = OpenAI(api_key=api_key)
        self.model = DEFAULT_MODEL
        self.full_content = ""
        self.temperature = 0.2
        logger.info("SlideGenerator initialized")
    
    def set_source_content(self, content: str):
        """Set source document content for context."""
        if not content or len(content.strip()) < 100:
            logger.error("Source content is empty or too short!")
            raise ValueError("Source document must contain at least 100 characters")
        
        self.full_content = content
        logger.info(f"Source content set: {len(self.full_content)} chars")
    
    @retry(wait=wait_exponential(multiplier=1, min=4, max=60), stop=stop_after_attempt(5))
    def generate_slides_for_module(
        self,
        module: Dict,
        progress_callback=None
    ) -> List[Dict]:
        """Generate slides for a module."""
        try:
            target_slides = min(module.get('estimated_slides', 8), 100)
            logger.info(f"Generating {target_slides} slides for module: {module.get('title')}")
            
            if progress_callback:
                progress_callback(f"Generating {target_slides} slides...", 0)
            
            context_chunks = chunk_text(self.full_content, 60000)
            base_context = "\n\n".join(context_chunks[:2])
            
            batch_size = 15
            all_slides = []
            
            for batch_start in range(0, target_slides, batch_size):
                batch_count = min(batch_size, target_slides - batch_start)
                batch_num = (batch_start // batch_size) + 1
                total_batches = (target_slides + batch_size - 1) // batch_size
                
                if progress_callback:
                    progress = int((batch_start / target_slides) * 100)
                    progress_callback(f"Batch {batch_num}/{total_batches}...", progress)
                
                previous_context = json.dumps(all_slides[-5:], default=str) if all_slides else "None"
                
                prompt = f"""Generate {batch_count} training slides for this module.

MODULE: {module.get('title')}
OBJECTIVES: {', '.join(module.get('objectives', [])[:3])}
TOPICS: {', '.join(module.get('topics', [])[:3])}

SOURCE CONTENT:
{base_context[:10000]}

PREVIOUS SLIDES:
{previous_context}

Return JSON array with {batch_count} slides:
[
  {{
    "slide_number": {batch_start + 1},
    "title": "Specific Topic Title",
    "content": [
      "Specific fact with data/numbers",
      "Another distinct detail",
      "Different concrete example",
      "Unique application"
    ],
    "notes": "Detailed 200+ word teaching guide",
    "slide_type": "content"
  }}
]

CRITICAL: Extract REAL information from source content above."""

                try:
                    response = self.client.chat.completions.create(
                        model=self.model,
                        messages=[
                            {
                                "role": "system",
                                "content": "Extract specific information from source documents. Return only valid JSON."
                            },
                            {
                                "role": "user",
                                "content": prompt
                            }
                        ],
                        response_format={"type": "json_object"},
                        temperature=0.1,
                        max_tokens=4000
                    )
                    
                    result = response.choices[0].message.content.strip()
                    parsed = safe_json_load(result)
                    
                    if isinstance(parsed, dict):
                        batch = parsed.get('slides', [])
                    elif isinstance(parsed, list):
                        batch = parsed
                    else:
                        continue
                    
                    for i, slide in enumerate(batch):
                        slide['slide_number'] = batch_start + i + 1
                    
                    all_slides.extend(batch)
                    
                except Exception as batch_error:
                    logger.warning(f"Batch {batch_num} failed: {str(batch_error)}")
                    continue
            
            final_slides = all_slides[:target_slides]
            
            for i, slide in enumerate(final_slides):
                if i == 0:
                    slide['slide_type'] = 'title'
                elif i == len(final_slides) - 1:
                    slide['slide_type'] = 'summary'
                else:
                    slide['slide_type'] = 'content'
            
            if progress_callback:
                progress_callback(f"Generated {len(final_slides)} slides", 100)
            
            return final_slides
            
        except Exception as e:
            logger.error(f"Slide generation error: {str(e)}")
            raise

class DocumentBuilder:
    """Creates final output documents with AI-powered assessments."""
    
    @staticmethod
    def create_powerpoint(
        outline: Dict,
        slides: List[Dict],
        output_path: Path,
        template: str = "Corporate"
    ) -> Path:
        """Create PowerPoint presentation."""
        try:
            logger.info(f"Creating PowerPoint with {len(slides)} slides")
            
            prs = PptxPresentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            
            title.text = outline.get('title', 'Training Program')
            subtitle.text = f"{outline.get('duration', 'N/A')} | {len(slides)} Slides"
            
            # Add all slides
            for slide_data in slides:
                if slide_data.get('slide_type') == 'title':
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    title = slide.shapes.title
                    title.text = slide_data.get('title', '')
                else:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    
                    title.text = slide_data.get('title', '')
                    
                    text_frame = content.text_frame
                    text_frame.clear()
                    
                    for bullet in slide_data.get('content', []):
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        p.space_before = Pt(6)
                
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = slide_data.get('notes', '')
            
            prs.save(str(output_path))
            logger.info(f"PowerPoint saved to {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"PowerPoint creation error: {str(e)}")
            raise
    
    @staticmethod
    def create_trainer_guide(
        outline: Dict,
        slides: List[Dict],
        output_path: Path
    ) -> Path:
        """Create comprehensive trainer guide document."""
        try:
            logger.info("Creating trainer guide")
            
            doc = DocxDocument()
            
            # Title page
            title = doc.add_heading(outline.get('title', 'Training Program'), 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            subtitle = doc.add_paragraph('Trainer Guide')
            subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_page_break()
            
            # Table of contents
            doc.add_heading('Table of Contents', 1)
            for i, module in enumerate(outline.get('modules', []), 1):
                doc.add_paragraph(f"Module {i}: {module.get('title')}", style='List Number')
            
            doc.add_page_break()
            
            # Program overview
            doc.add_heading('Program Overview', 1)
            doc.add_paragraph(outline.get('description', ''))
            
            doc.add_heading('Duration', 2)
            doc.add_paragraph(outline.get('duration', 'N/A'))
            
            doc.add_heading('Learning Objectives', 2)
            for obj in outline.get('objectives', []):
                doc.add_paragraph(obj, style='List Bullet')
            
            doc.add_page_break()
            
            # Module details
            slide_index = 0
            for module_num, module in enumerate(outline.get('modules', []), 1):
                doc.add_heading(f"Module {module_num}: {module.get('title')}", 1)
                
                doc.add_heading('Module Overview', 2)
                doc.add_paragraph(f"Duration: {module.get('duration')}")
                doc.add_paragraph(f"Estimated Slides: {module.get('estimated_slides')}")
                
                doc.add_heading('Module Objectives', 2)
                for obj in module.get('objectives', []):
                    doc.add_paragraph(obj, style='List Bullet')
                
                doc.add_heading('Key Points', 2)
                for kp in module.get('key_points', []):
                    doc.add_paragraph(kp, style='List Bullet')
                
                doc.add_heading('Slide-by-Slide Guide', 2)
                
                module_slides = module.get('estimated_slides', 10)
                for i in range(module_slides):
                    if slide_index < len(slides):
                        slide = slides[slide_index]
                        
                        doc.add_heading(f"Slide {slide_index + 1}: {slide.get('title')}", 3)
                        
                        doc.add_paragraph('Content:')
                        for bullet in slide.get('content', []):
                            doc.add_paragraph(bullet, style='List Bullet 2')
                        
                        doc.add_paragraph('Trainer Notes:')
                        doc.add_paragraph(slide.get('notes', ''))
                        
                        slide_index += 1
                
                doc.add_page_break()
            
            doc.save(str(output_path))
            logger.info(f"Trainer guide saved to {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"Trainer guide creation error: {str(e)}")
            raise
    
    @staticmethod
    @retry(wait=wait_exponential(multiplier=1, min=4, max=30), stop=stop_after_attempt(3))
    def _generate_assessment_questions(
        api_key: str,
        module: Dict,
        source_content: str
    ) -> List[Dict]:
        """AI-powered assessment question generation."""
        try:
            client = OpenAI(api_key=api_key)
            
            context_chunk = source_content[:5000] if source_content else ""
            
            prompt = f"""Generate assessment questions for this training module.

MODULE: {module.get('title')}
OBJECTIVES: {', '.join(module.get('objectives', [])[:3])}
KEY POINTS: {', '.join(module.get('key_points', [])[:5])}

SOURCE CONTENT:
{context_chunk}

Generate 4 questions in JSON format:
{{
  "questions": [
    {{
      "type": "multiple_choice",
      "question": "Specific question based on content",
      "options": ["A. Option 1", "B. Option 2", "C. Option 3", "D. Option 4"],
      "correct_answer": "A",
      "explanation": "Why this is correct"
    }},
    {{
      "type": "multiple_choice",
      "question": "Another specific question",
      "options": ["A. ...", "B. ...", "C. ...", "D. ..."],
      "correct_answer": "B",
      "explanation": "Explanation"
    }},
    {{
      "type": "multiple_choice",
      "question": "Third question",
      "options": ["A. ...", "B. ...", "C. ...", "D. ..."],
      "correct_answer": "C",
      "explanation": "Explanation"
    }},
    {{
      "type": "short_answer",
      "question": "Open-ended question about key concept",
      "grading_points": ["Point 1", "Point 2", "Point 3"],
      "sample_answer": "Example good answer"
    }}
  ]
}}

Requirements:
- Use REAL content from source document
- Include specific data, numbers, or facts
- Make questions practical and relevant
- Ensure correct answers are accurate
- Return ONLY valid JSON"""

            response = client.chat.completions.create(
                model=DEFAULT_MODEL,
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert assessment designer. Create specific, relevant questions based on source content. Return only valid JSON."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                response_format={"type": "json_object"},
                temperature=0.3,
                max_tokens=2000
            )
            
            result = safe_json_load(response.choices[0].message.content)
            questions = result.get('questions', [])
            
            logger.info(f"Generated {len(questions)} questions for module: {module.get('title')}")
            return questions
            
        except Exception as e:
            logger.warning(f"AI question generation failed: {str(e)}")
            return DocumentBuilder._generate_fallback_questions(module)
    
    @staticmethod
    def _generate_fallback_questions(module: Dict) -> List[Dict]:
        """Fallback: Generate generic questions from module objectives."""
        questions = []
        
        for i, objective in enumerate(module.get('objectives', [])[:3], 1):
            questions.append({
                'type': 'multiple_choice',
                'question': f"Which statement best describes: {objective}?",
                'options': [
                    "A. First interpretation",
                    "B. Second interpretation",
                    "C. Third interpretation",
                    "D. Fourth interpretation"
                ],
                'correct_answer': 'A',
                'explanation': f"Based on module objective: {objective}"
            })
        
        key_points = module.get('key_points', [])
        if key_points:
            questions.append({
                'type': 'short_answer',
                'question': f"Explain the key concepts covered in {module.get('title')}.",
                'grading_points': key_points[:3],
                'sample_answer': f"A good answer should cover: {', '.join(key_points[:3])}"
            })
        
        return questions
    
    @staticmethod
    def create_assessment(
        outline: Dict,
        output_path: Path,
        api_key: str = None,
        source_content: str = None
    ) -> Path:
        """Create AI-powered assessment document with real questions."""
        try:
            logger.info("Creating AI-powered assessment document")
            
            doc = DocxDocument()
            
            # Title
            title = doc.add_heading(f"{outline.get('title')} - Assessment", 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_paragraph(f"Duration: {outline.get('duration')}")
            doc.add_paragraph(f"Total Modules: {outline.get('total_modules')}")
            doc.add_paragraph("")
            
            doc.add_heading("Instructions", 2)
            doc.add_paragraph("Please answer all questions to the best of your ability.")
            doc.add_paragraph("For multiple choice questions, select the best answer.")
            doc.add_paragraph("For short answer questions, provide detailed responses.")
            
            doc.add_page_break()
            
            # Generate questions for each module
            question_num = 1
            
            for module in outline.get('modules', []):
                doc.add_heading(f"Module: {module.get('title')}", 1)
                doc.add_paragraph("")
                
                # Generate AI-powered questions
                if api_key and source_content:
                    try:
                        questions = DocumentBuilder._generate_assessment_questions(
                            api_key, module, source_content
                        )
                    except Exception:
                        questions = DocumentBuilder._generate_fallback_questions(module)
                else:
                    questions = DocumentBuilder._generate_fallback_questions(module)
                
                # Add questions to document
                for q in questions:
                    if q['type'] == 'multiple_choice':
                        doc.add_heading(f"Question {question_num}", 2)
                        doc.add_paragraph(q['question'])
                        doc.add_paragraph("")
                        
                        for option in q['options']:
                            if option.startswith(q['correct_answer']):
                                p = doc.add_paragraph(f"{option} ✓", style='List Bullet')
                                for run in p.runs:
                                    run.bold = True
                            else:
                                doc.add_paragraph(option, style='List Bullet')
                        
                        doc.add_paragraph("")
                        doc.add_paragraph(f"Correct Answer: {q['correct_answer']}")
                        doc.add_paragraph(f"Explanation: {q.get('explanation', '')}")
                        doc.add_paragraph("")
                        
                    elif q['type'] == 'short_answer':
                        doc.add_heading(f"Question {question_num}", 2)
                        doc.add_paragraph(q['question'])
                        doc.add_paragraph("")
                        doc.add_paragraph("Answer:")
                        doc.add_paragraph("_" * 80)
                        doc.add_paragraph("_" * 80)
                        doc.add_paragraph("_" * 80)
                        doc.add_paragraph("")
                        
                        doc.add_paragraph("Grading Points:")
                        for point in q.get('grading_points', []):
                            doc.add_paragraph(f"• {point}", style='List Bullet 2')
                        doc.add_paragraph("")
                        
                        doc.add_paragraph(f"Sample Answer: {q.get('sample_answer', '')}")
                        doc.add_paragraph("")
                    
                    question_num += 1
                
                doc.add_page_break()
            
            # Answer key summary
            doc.add_heading("Answer Key Summary", 1)
            doc.add_paragraph("Quick reference for all multiple choice answers:")
            doc.add_paragraph("")
            
            q_num = 1
            for module in outline.get('modules', []):
                if api_key and source_content:
                    try:
                        questions = DocumentBuilder._generate_assessment_questions(
                            api_key, module, source_content
                        )
                    except Exception:
                        questions = DocumentBuilder._generate_fallback_questions(module)
                else:
                    questions = DocumentBuilder._generate_fallback_questions(module)
                
                for q in questions:
                    if q['type'] == 'multiple_choice':
                        doc.add_paragraph(f"Q{q_num}: {q['correct_answer']}")
                        q_num += 1
            
            doc.save(str(output_path))
            logger.info(f"AI-powered assessment saved to {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"Assessment creation error: {str(e)}")
            raise
    
    @staticmethod
    def create_zip_package(file_paths: List[Path], output_path: Path) -> Path:
        """Create ZIP package of all documents."""
        try:
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for file_path in file_paths:
                    if file_path.exists():
                        zipf.write(file_path, file_path.name)
            
            logger.info(f"ZIP package created: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"ZIP creation error: {str(e)}")
            raise

# ============================================================================
# CONTENT EDITOR CLASS WITH IMPROVED IMPORTANCE DISPLAY
# ============================================================================

class ContentEditor:
    """Provides interactive editing interfaces for topics and outlines."""
    
    @staticmethod
    def render_topics_editor(topics: List[Dict]) -> List[Dict]:
        """Render interactive UI for editing topics WITH DELETE and ADD functionality."""
        st.markdown("### 📋 Edit Extracted Topics")
        st.markdown("Review and modify the topics extracted from your document.")
        
        # Initialize delete tracking
        if 'topics_to_delete' not in st.session_state:
            st.session_state.topics_to_delete = []
        
        edited_topics = []
        
        # Filter out deleted topics
        active_topics = [t for t in topics if t.get('id') not in st.session_state.topics_to_delete]
        
        for idx, topic in enumerate(active_topics):
            topic_id = topic.get('id', idx + 1)
            importance_badge = render_importance_badge(topic.get('importance', 'medium'))
            
            with st.expander(
                f"**Topic {idx + 1}: {topic.get('title', 'Untitled')}**", 
                expanded=(idx == 0)
            ):
                # Delete button at the top
                col_header = st.columns([4, 1])
                with col_header[0]:
                    st.markdown(f"Current Importance: {importance_badge}", unsafe_allow_html=True)
                with col_header[1]:
                    if st.button("🗑️ Delete", key=f"delete_topic_{topic_id}", type="secondary"):
                        st.session_state.topics_to_delete.append(topic_id)
                        st.rerun()
                
                st.markdown("---")
                
                col1, col2 = st.columns([3, 1])
                
                with col1:
                    title = st.text_input(
                        "Topic Title",
                        value=topic.get('title', ''),
                        key=f"topic_title_{topic_id}"
                    )
                
                with col2:
                    # Color-coded importance selector
                    current_importance = topic.get('importance', 'medium')
                    importance_options = {
                        'high': '🔴 HIGH (Critical)',
                        'medium': '🟡 MEDIUM (Important)',
                        'low': '🔵 LOW (Optional)'
                    }
                    
                    importance_display = st.selectbox(
                        "Importance Level",
                        options=list(importance_options.values()),
                        index=list(importance_options.keys()).index(current_importance),
                        key=f"topic_importance_{topic_id}"
                    )
                    
                    # Extract actual importance value
                    importance = list(importance_options.keys())[
                        list(importance_options.values()).index(importance_display)
                    ]
                
                description = st.text_area(
                    "Description",
                    value=topic.get('description', ''),
                    key=f"topic_desc_{topic_id}",
                    height=100
                )
                
                col3, col4 = st.columns(2)
                
                with col3:
                    key_concepts_str = ', '.join(topic.get('key_concepts', []))
                    key_concepts = st.text_input(
                        "Key Concepts (comma-separated)",
                        value=key_concepts_str,
                        key=f"topic_concepts_{topic_id}"
                    )
                
                with col4:
                    duration = st.number_input(
                        "Duration (minutes)",
                        min_value=10,
                        max_value=180,
                        value=int(topic.get('estimated_duration_minutes', 45)),
                        step=5,
                        key=f"topic_duration_{topic_id}"
                    )
                
                edited_topics.append({
                    'id': topic_id,
                    'title': title,
                    'description': description,
                    'key_concepts': [k.strip() for k in key_concepts.split(',') if k.strip()],
                    'importance': importance,
                    'estimated_duration_minutes': duration
                })
        
        st.markdown("---")
        
        # Action buttons - ALL WITH UNIQUE KEYS
        col1, col2, col3 = st.columns([1, 2, 1])
        
        with col1:
            if st.button("➕ Add New Topic", use_container_width=True, key="add_new_topic_btn"):
                # Generate new unique ID
                existing_ids = [t.get('id', 0) for t in edited_topics]
                new_id = max(existing_ids + [0]) + 1
                
                new_topic = {
                    'id': new_id,
                    'title': f"New Topic {len(edited_topics) + 1}",
                    'description': "Enter topic description here",
                    'key_concepts': ["concept1", "concept2", "concept3"],
                    'importance': 'medium',
                    'estimated_duration_minutes': 45
                }
                edited_topics.append(new_topic)
                
                # Save immediately to session state
                st.session_state.edited_topics = edited_topics
                st.success("✅ New topic added!")
                time.sleep(0.5)
                st.rerun()
        
        with col2:
            if st.button("💾 Save Changes", use_container_width=True, type="primary", key="save_topics_btn"):
                st.session_state.edited_topics = edited_topics
                st.session_state.topics_to_delete = []  # Clear delete tracking
                st.success("✅ Topics saved successfully!")
                st.balloons()
                time.sleep(1)
                st.rerun()
        
        with col3:
            if st.button("🔄 Reset", use_container_width=True, key="reset_topics_btn"):
                st.session_state.topics_to_delete = []
                st.session_state.edited_topics = None
                st.rerun()
        
        # Show current stats
        st.info(f"📊 Current topics: {len(edited_topics)} topics, Total duration: {sum(t['estimated_duration_minutes'] for t in edited_topics)} minutes")
        
        return edited_topics
    
    @staticmethod
    def render_outline_editor(outline: Dict) -> Dict:
        """Render interactive UI for editing training outline with delete functionality."""
        st.markdown("### 📝 Edit Training Outline")
        st.markdown("Customize your training program structure and modules.")
        
        # Initialize delete tracking
        if 'modules_to_delete' not in st.session_state:
            st.session_state.modules_to_delete = []
        
        st.markdown("#### Global Settings")
        
        col1, col2 = st.columns(2)
        
        with col1:
            title = st.text_input(
                "Training Program Title",
                value=outline.get('title', ''),
                key="outline_title"
            )
        
        with col2:
            duration = st.selectbox(
                "Training Duration",
                options=list(DURATION_TO_SLIDES.keys()),
                index=list(DURATION_TO_SLIDES.keys()).index(outline.get('duration', '1 day')) 
                    if outline.get('duration') in DURATION_TO_SLIDES else 4,
                key="outline_duration"
            )
        
        description = st.text_area(
            "Program Description",
            value=outline.get('description', ''),
            key="outline_description",
            height=100
        )
        
        objectives_text = '\n'.join(outline.get('objectives', []))
        objectives = st.text_area(
            "Learning Objectives (one per line)",
            value=objectives_text,
            key="outline_objectives",
            height=120
        )
        objectives_list = [obj.strip() for obj in objectives.split('\n') if obj.strip()]
        
        st.markdown("---")
        st.markdown("#### Training Modules")
        
        edited_modules = []
        modules_list = outline.get('modules', [])
        
        # Filter out deleted modules
        modules_list = [m for m in modules_list if m.get('id') not in st.session_state.modules_to_delete]
        
        for idx, module in enumerate(modules_list):
            module_id = module.get('id', idx + 1)
            
            with st.expander(
                f"**Module {idx + 1}: {module.get('title', 'Untitled')}** "
                f"({module.get('estimated_slides', 0)} slides)",
                expanded=(idx == 0)
            ):
                # Delete button at the top of each module - UNIQUE KEY
                col_delete = st.columns([4, 1])
                with col_delete[0]:
                    st.markdown(f"**Module ID: {module_id}**")
                with col_delete[1]:
                    if st.button(f"🗑️ Delete", key=f"delete_module_{module_id}", type="secondary"):
                        st.session_state.modules_to_delete.append(module_id)
                        st.success(f"Module {idx + 1} deleted")
                        time.sleep(0.5)
                        st.rerun()
                
                col1, col2, col3 = st.columns([2, 1, 1])
                
                with col1:
                    module_title = st.text_input(
                        "Module Title",
                        value=module.get('title', ''),
                        key=f"module_title_{module_id}"
                    )
                
                with col2:
                    module_duration = st.text_input(
                        "Duration",
                        value=module.get('duration', ''),
                        key=f"module_duration_{module_id}"
                    )
                
                with col3:
                    estimated_slides = st.number_input(
                        "Slides",
                        min_value=1,
                        max_value=100,
                        value=min(int(module.get('estimated_slides', 10)), 100),
                        key=f"module_slides_{module_id}",
                        help="Use +/- buttons to adjust"
                    )
                
                module_objectives_text = '\n'.join(module.get('objectives', []))
                module_objectives = st.text_area(
                    "Module Objectives (one per line)",
                    value=module_objectives_text,
                    key=f"module_objectives_{module_id}",
                    height=100
                )
                module_objectives_list = [obj.strip() for obj in module_objectives.split('\n') if obj.strip()]
                
                topics_covered_str = ', '.join(module.get('topics', []))
                topics_covered = st.text_input(
                    "Topics Covered (comma-separated)",
                    value=topics_covered_str,
                    key=f"module_topics_{module_id}"
                )
                topics_covered_list = [t.strip() for t in topics_covered.split(',') if t.strip()]
                
                key_points_text = '\n'.join(module.get('key_points', []))
                key_points = st.text_area(
                    "Key Points (one per line)",
                    value=key_points_text,
                    key=f"module_keypoints_{module_id}",
                    height=120
                )
                key_points_list = [kp.strip() for kp in key_points.split('\n') if kp.strip()]
                
                edited_modules.append({
                    'id': module_id,
                    'title': module_title,
                    'duration': module_duration,
                    'objectives': module_objectives_list,
                    'topics': topics_covered_list,
                    'key_points': key_points_list,
                    'estimated_slides': estimated_slides
                })
        
        st.markdown("---")
        
        # Action buttons - ALL WITH UNIQUE KEYS
        col1, col2, col3 = st.columns([1, 2, 1])
        
        with col1:
            if st.button("➕ Add Module", use_container_width=True, key="add_new_module_btn"):
                # Generate new unique ID
                existing_ids = [m.get('id', 0) for m in edited_modules]
                new_id = max(existing_ids + [0]) + 1
                
                new_module = {
                    'id': new_id,
                    'title': f"New Module {len(edited_modules) + 1}",
                    'duration': "2 hours",
                    'objectives': ["Objective 1", "Objective 2", "Objective 3"],
                    'topics': ["Topic 1", "Topic 2"],
                    'key_points': ["Key point 1", "Key point 2", "Key point 3"],
                    'estimated_slides': 10
                }
                edited_modules.append(new_module)
                
                # Update the outline in session state immediately
                edited_outline = {
                    'title': title,
                    'description': description,
                    'duration': duration,
                    'objectives': objectives_list,
                    'modules': edited_modules,
                    'total_modules': len(edited_modules),
                    'estimated_slides': sum(m['estimated_slides'] for m in edited_modules)
                }
                st.session_state.edited_outline = edited_outline
                st.success("✅ New module added!")
                time.sleep(0.5)
                st.rerun()
        
        with col2:
            if st.button("💾 Save Changes", use_container_width=True, type="primary", key="save_outline_btn"):
                edited_outline = {
                    'title': title,
                    'description': description,
                    'duration': duration,
                    'objectives': objectives_list,
                    'modules': edited_modules,
                    'total_modules': len(edited_modules),
                    'estimated_slides': sum(m['estimated_slides'] for m in edited_modules)
                }
                st.session_state.edited_outline = edited_outline
                st.session_state.modules_to_delete = []  # Clear delete tracking
                st.success("✅ Changes saved successfully!")
                st.balloons()
                time.sleep(1)
                st.rerun()
        
        with col3:
            if st.button("🔄 Reset", use_container_width=True, key="reset_outline_btn"):
                st.session_state.modules_to_delete = []
                st.session_state.edited_outline = None
                st.rerun()
        
        # Show current stats
        st.info(f"📊 Current outline: {len(edited_modules)} modules, {sum(m['estimated_slides'] for m in edited_modules)} total slides")
        
        edited_outline = {
            'title': title,
            'description': description,
            'duration': duration,
            'objectives': objectives_list,
            'modules': edited_modules,
            'total_modules': len(edited_modules),
            'estimated_slides': sum(m['estimated_slides'] for m in edited_modules)
        }
        
        return edited_outline

# ============================================================================
# PAGE SETUP WITH THEME TOGGLE
# ============================================================================

def setup_page():
    """Configure Streamlit page and apply custom styling."""
    st.set_page_config(
        page_title="Enhanced Training Generator",
        page_icon="📚",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    initialize_session_state()
    
    # Apply theme-specific CSS
    theme_css = get_theme_css(st.session_state.theme)
    st.markdown(theme_css, unsafe_allow_html=True)

def display_header():
    """Display main application header with theme toggle."""
    col1, col2, col3 = st.columns([1, 3, 1])
    
    with col1:
        # Theme toggle button
        theme_icon = "🌙" if st.session_state.theme == "dark" else "☀️"
        theme_text = "Light Mode" if st.session_state.theme == "dark" else "Dark Mode"
        
        if st.button(f"{theme_icon} {theme_text}", key="theme_toggle"):
            st.session_state.theme = "light" if st.session_state.theme == "dark" else "dark"
            st.rerun()
    
    with col2:
        st.markdown("""
        <div style="text-align: center; padding: 2rem 0;">
            <div style="display: flex; align-items: center; justify-content: center; margin-bottom: 1rem;">
                <span style="font-size: 3rem;">✨</span>
                <h1 style="font-size: 3.5rem; margin: 0 1.5rem; background: linear-gradient(135deg, #667eea 0%, #764ba2 50%, #f093fb 100%); -webkit-background-clip: text; -webkit-text-fill-color: transparent; font-weight: 800;">
                    QTS Training Generator
                </h1>
                <span style="font-size: 3rem;">💫</span>
            </div>
            <p style="font-size: 1.3rem; margin-bottom: 1rem;">
                Transform documents into professional training materials
            </p>
            
        </div>
        """, unsafe_allow_html=True)
    
    
def display_phase_tracker():
    """Display phase tracker."""
    phases = [
        ("📄", "Upload", PHASE_1, "#667eea"),
        ("🔍", "Analyze", PHASE_2, "#764ba2"),
        ("✏️", "Edit", PHASE_3, "#f093fb"),
        ("🎨", "Generate", PHASE_4, "#4facfe")
    ]
    
    st.markdown("<div style='padding: 2rem 0;'>", unsafe_allow_html=True)
    
    cols = st.columns(4)
    
    for col, (icon, name, phase_id, color) in zip(cols, phases):
        with col:
            is_current = st.session_state.current_phase == phase_id
            is_complete = st.session_state.phase_completed.get(phase_id, False)
            
            if is_complete:
                badge_class = "phase-complete"
                display_icon = "✅"
            elif is_current:
                badge_class = "phase-active"
                display_icon = icon
            else:
                badge_class = "phase-pending"
                display_icon = icon
            
            st.markdown(f"""
            <div style="text-align: center;">
                <div style="width: 70px; height: 70px; margin: 0 auto 0.5rem; border-radius: 50%; 
                     display: flex; align-items: center; justify-content: center; font-size: 2rem;
                     background: {'linear-gradient(135deg, ' + color + ', ' + color + 'dd)' if (is_current or is_complete) else 'rgba(128, 128, 128, 0.3)'};
                     border: 2px solid {'rgba(255, 255, 255, 0.5)' if (is_current or is_complete) else 'rgba(128, 128, 128, 0.3)'};
                     box-shadow: {'0 4px 15px ' + color + '40' if is_current else 'none'};
                     transition: all 0.3s ease;">
                    {display_icon}
                </div>
                <div class="phase-badge {badge_class}" style="width: 100%; font-size: 0.9rem;">
                    {name}
                </div>
            </div>
            """, unsafe_allow_html=True)
    
    st.markdown("</div>", unsafe_allow_html=True)

# ============================================================================
# SIDEBAR WITH IMPROVED API KEY DISPLAY
# ============================================================================

def render_sidebar():
    """Render sidebar with configuration options."""
    with st.sidebar:
        st.markdown("## ⚙️ Configuration")
        
        # API Key section with auto-load info
        st.markdown("### 🔑 API Key")
        
        if st.session_state.api_key:
            st.success(f"✅ API Key loaded")
            st.caption("Key loaded from secrets.toml")
            
            if st.button("🔄 Use Different Key", use_container_width=True):
                manual_key = st.text_input(
                    "Enter New API Key",
                    type="password",
                    key="manual_api_key_input"
                )
                if manual_key and validate_api_key(manual_key):
                    st.session_state.api_key = manual_key
                    st.success("✅ New API key set!")
                    time.sleep(1)
                    st.rerun()
        else:
            st.warning("⚠️ No API key found")
            st.info("💡 Add OPENAI_API_KEY to .streamlit/secrets.toml file or enter manually below")
            
            api_key = st.text_input(
                "OpenAI API Key",
                type="password",
                key="api_key_input"
            )
            
            if api_key:
                if validate_api_key(api_key):
                    st.session_state.api_key = api_key
                    st.success("✅ API key validated")
                    time.sleep(1)
                    st.rerun()
                else:
                    st.error("❌ Invalid API key format")
        
        st.markdown("---")
        
        # Training Settings
        st.markdown("### 📊 Training Settings")
        
        duration = st.selectbox(
            "Training Duration",
            options=list(DURATION_TO_SLIDES.keys()),
            index=list(DURATION_TO_SLIDES.keys()).index(st.session_state.training_duration)
        )
        st.session_state.training_duration = duration
        
        recommended_slides = DURATION_TO_SLIDES[duration]
        
        use_custom = st.checkbox(
            "Custom Slide Count",
            value=st.session_state.custom_slide_count is not None
        )
        
        if use_custom:
            custom_slides = st.number_input(
                "Number of Slides",
                min_value=10,
                max_value=720,
                value=max(10, st.session_state.custom_slide_count or recommended_slides),
                step=10,
                help="Use the +/- buttons to adjust"
            )
            st.session_state.custom_slide_count = custom_slides
            st.session_state.target_slides = custom_slides
        else:
            st.session_state.custom_slide_count = None
            st.session_state.target_slides = recommended_slides
            st.info(f"📊 Recommended: {recommended_slides} slides")
        
        modules = st.slider(
            "Number of Modules",
            min_value=5,
            max_value=15,
            value=st.session_state.target_modules
        )
        st.session_state.target_modules = modules
        
        st.markdown("---")
        
        # Template
        st.markdown("### 🎨 Template Style")
        template = st.selectbox(
            "Select Template",
            options=TEMPLATE_OPTIONS,
            index=TEMPLATE_OPTIONS.index(st.session_state.template)
        )
        st.session_state.template = template
        
        st.markdown("---")
        
        # Enhancements
        st.markdown("### ✨ Enhancements")
        
        st.session_state.include_assessments = st.checkbox(
            "Include AI Assessments",
            value=st.session_state.include_assessments,
            help="Generate AI-powered questions with real content"
        )
        
        st.session_state.enhanced_trainer_guide = st.checkbox(
            "Enhanced Trainer Guide",
            value=st.session_state.enhanced_trainer_guide
        )
        
        st.markdown("---")
        
        # Progress
        st.markdown("### 📈 Progress")
        
        completed = sum(st.session_state.phase_completed.values())
        progress = completed / 4
        
        st.progress(progress)
        st.caption(f"{completed}/4 phases completed")
        
        st.markdown("---")
        
        # Quick Navigation
        st.markdown("### 🔄 Quick Navigation")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("📄 Phase 1", use_container_width=True):
                st.session_state.current_phase = PHASE_1
                st.rerun()
            if st.button("✏️ Phase 3", use_container_width=True, disabled=not st.session_state.phase_completed.get(PHASE_2)):
                st.session_state.current_phase = PHASE_3
                st.rerun()
        
        with col2:
            if st.button("🔍 Phase 2", use_container_width=True, disabled=not st.session_state.phase_completed.get(PHASE_1)):
                st.session_state.current_phase = PHASE_2
                st.rerun()
            if st.button("🎨 Phase 4", use_container_width=True, disabled=not st.session_state.phase_completed.get(PHASE_3)):
                st.session_state.current_phase = PHASE_4
                st.rerun()

# ============================================================================
# PHASE RENDERING (Same logic as before, just cleaner)
# ============================================================================

def render_phase_1():
    """Render Phase 1: Document Upload."""
    st.markdown("## 📄 Phase 1: Document Upload & Processing")
    
    if st.session_state.phase_completed.get(PHASE_1):
        st.success("✅ Phase 1 completed!")
        if st.button("➡️ Continue to Phase 2", type="primary", use_container_width=True):
            st.session_state.current_phase = PHASE_2
            st.rerun()
        st.markdown("---")
    
    uploaded_files = st.file_uploader(
        "Choose document file(s)",
        type=SUPPORTED_FORMATS,
        accept_multiple_files=True
    )
    
    if uploaded_files:
        st.success(f"✅ {len(uploaded_files)} file(s) uploaded")
        
        if st.button("🔄 Process Document(s)", type="primary", use_container_width=True):
            if not st.session_state.api_key:
                st.error("⚠️ Please set your OpenAI API key")
                return
            
            with st.spinner("Processing..."):
                try:
                    processor = DocumentProcessor()
                    combined_text = []
                    total_words = 0
                    
                    for file in uploaded_files:
                        processed = processor.process_file(file)
                        combined_text.append(processed['text'])
                        total_words += processed['word_count']
                    
                    final_text = '\n\n'.join(combined_text)
                    
                    st.session_state.full_source_text = final_text
                    st.session_state.processed_content = {
                        'text': final_text[:20000],
                        'word_count': total_words,
                        'file_count': len(uploaded_files)
                    }
                    st.session_state.phase_completed[PHASE_1] = True
                    
                    st.success(f"✅ Processed {total_words} words from {len(uploaded_files)} file(s)")
                    
                    if st.button("➡️ Proceed to Phase 2", type="primary", use_container_width=True):
                        st.session_state.current_phase = PHASE_2
                        st.rerun()
                    
                except Exception as e:
                    st.error(f"❌ Error: {str(e)}")
    else:
        st.info("📁 Upload documents to begin")

def render_phase_2():
    """Render Phase 2: Topic Analysis."""
    st.markdown("## 🔍 Phase 2: Topic Analysis")
    
    if not st.session_state.processed_content:
        st.warning("⚠️ Complete Phase 1 first")
        return
    
    if st.button("🚀 Extract Topics & Generate Outline", type="primary", use_container_width=True):
        if not st.session_state.api_key:
            st.error("⚠️ Enter API key first")
            return
        
        with st.spinner("Analyzing..."):
            try:
                analyzer = TopicAnalyzer(st.session_state.api_key)
                topics = analyzer.extract_topics(
                    st.session_state.full_source_text,
                    num_topics=st.session_state.target_modules
                )
                
                generator = OutlineGenerator(st.session_state.api_key)
                outline = generator.generate_outline(
                    topics=topics,
                    target_modules=st.session_state.target_modules,
                    target_slides=st.session_state.target_slides,
                    duration=st.session_state.training_duration
                )
                
                st.session_state.extracted_topics = topics
                st.session_state.generated_outline = outline
                st.session_state.phase_completed[PHASE_2] = True
                
                st.success("✅ Analysis complete!")
                
                if st.button("➡️ Proceed to Phase 3", type="primary", use_container_width=True):
                    st.session_state.current_phase = PHASE_3
                    st.rerun()
                
            except Exception as e:
                st.error(f"❌ Error: {str(e)}")
    
    elif st.session_state.generated_outline:
        st.info("✅ Analysis completed")
        if st.button("➡️ Proceed to Phase 3", type="primary", use_container_width=True):
            st.session_state.current_phase = PHASE_3
            st.rerun()

def render_phase_3():
    """Render Phase 3: Content Editing."""
    st.markdown("## ✏️ Phase 3: Content Review & Editing")
    
    if not st.session_state.generated_outline:
        st.warning("⚠️ Complete Phase 2 first")
        return
    
    tab1, tab2 = st.tabs(["📋 Topics", "📝 Outline"])
    
    with tab1:
        topics = st.session_state.edited_topics or st.session_state.extracted_topics
        edited_topics = ContentEditor.render_topics_editor(topics)
        st.session_state.edited_topics = edited_topics
    
    with tab2:
        outline = st.session_state.edited_outline or st.session_state.generated_outline
        edited_outline = ContentEditor.render_outline_editor(outline)
        # Note: edited_outline is returned but saved via button click
    
    st.markdown("---")
    
    # Only allow proceeding if changes have been saved
    if st.session_state.edited_outline:
        if st.button("✅ Finalize & Continue to Phase 4", type="primary", use_container_width=True):
            st.session_state.phase_completed[PHASE_3] = True
            st.session_state.current_phase = PHASE_4
            st.rerun()
    else:
        st.info("💡 Click 'Save Changes' in the Outline tab to proceed to Phase 4")

def render_phase_4():
    """Render Phase 4: Document Generation."""
    st.markdown("## 🎨 Phase 4: Final Document Generation")
    
    if not st.session_state.edited_outline:
        st.warning("⚠️ Complete Phase 3 first")
        return
    
    outline = st.session_state.edited_outline
    
    st.markdown("### 📊 Generation Summary")
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("Modules", outline['total_modules'])
    with col2:
        st.metric("Slides", outline['estimated_slides'])
    with col3:
        st.metric("Duration", outline['duration'])
    with col4:
        st.metric("Assessment", "AI-Powered" if st.session_state.include_assessments else "None")
    
    st.markdown("---")
    
    if st.button("🚀 Generate Final Documents", type="primary", use_container_width=True):
        if not st.session_state.api_key:
            st.error("⚠️ Enter API key first")
            return
        
        with st.spinner("Generating documents..."):
            try:
                # Generate slides
                slide_gen = SlideGenerator(st.session_state.api_key)
                slide_gen.set_source_content(st.session_state.full_source_text)
                
                all_slides = []
                for module in outline['modules']:
                    slides = slide_gen.generate_slides_for_module(module)
                    all_slides.extend(slides)
                
                st.session_state.generated_slides = all_slides
                
                # Create documents
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                generated_files = {}
                
                # PowerPoint
                pptx_path = OUTPUT_DIR / f"training_{timestamp}.pptx"
                DocumentBuilder.create_powerpoint(outline, all_slides, pptx_path)
                generated_files['powerpoint'] = pptx_path
                
                # Trainer Guide
                trainer_path = OUTPUT_DIR / f"trainer_guide_{timestamp}.docx"
                DocumentBuilder.create_trainer_guide(outline, all_slides, trainer_path)
                generated_files['trainer_guide'] = trainer_path
                
                # AI-Powered Assessment
                if st.session_state.include_assessments:
                    assessment_path = OUTPUT_DIR / f"assessment_{timestamp}.docx"
                    DocumentBuilder.create_assessment(
                        outline,
                        assessment_path,
                        api_key=st.session_state.api_key,
                        source_content=st.session_state.full_source_text
                    )
                    generated_files['assessment'] = assessment_path
                
                # ZIP Package
                zip_path = OUTPUT_DIR / f"training_package_{timestamp}.zip"
                DocumentBuilder.create_zip_package(list(generated_files.values()), zip_path)
                generated_files['zip_package'] = zip_path
                
                st.session_state.final_documents = generated_files
                st.session_state.phase_completed[PHASE_4] = True
                
                st.success(f"🎉 Generated {len(all_slides)} slides across {len(outline['modules'])} modules!")
                
                st.markdown("### 📥 Download Your Materials")
                
                cols = st.columns(len(generated_files))
                icons = {'powerpoint': '📊', 'trainer_guide': '📖', 'assessment': '📋', 'zip_package': '📦'}
                
                for col, (doc_type, file_path) in zip(cols, generated_files.items()):
                    with col:
                        with open(file_path, 'rb') as f:
                            st.download_button(
                                label=f"{icons.get(doc_type, '📄')} {doc_type.replace('_', ' ').title()}",
                                data=f.read(),
                                file_name=file_path.name,
                                use_container_width=True
                            )
                
            except Exception as e:
                st.error(f"❌ Error: {str(e)}")
    
    elif st.session_state.final_documents:
        st.success("✅ Documents already generated!")
        
        st.markdown("### 📥 Download Your Materials")
        
        cols = st.columns(len(st.session_state.final_documents))
        icons = {'powerpoint': '📊', 'trainer_guide': '📖', 'assessment': '📋', 'zip_package': '📦'}
        
        for col, (doc_type, file_path) in zip(cols, st.session_state.final_documents.items()):
            with col:
                if file_path.exists():
                    with open(file_path, 'rb') as f:
                        st.download_button(
                            label=f"{icons.get(doc_type, '📄')} {doc_type.replace('_', ' ').title()}",
                            data=f.read(),
                            file_name=file_path.name,
                            use_container_width=True
                        )

# ============================================================================
# MAIN APP
# ============================================================================

def main():
    """Main application entry point."""
    setup_page()
    display_header()
    display_phase_tracker()
    
    st.markdown("---")
    
    render_sidebar()
    
    current_phase = st.session_state.current_phase
    
    if current_phase == PHASE_1:
        render_phase_1()
    elif current_phase == PHASE_2:
        render_phase_2()
    elif current_phase == PHASE_3:
        render_phase_3()
    elif current_phase == PHASE_4:
        render_phase_4()
    
    st.markdown("---")
    st.markdown(f"""
    <div style="text-align: center; padding: 2rem 0; opacity: 0.8;">
        <p>Enhanced Training Generator  | {ORGANIZATION}</p>
        <p style="font-size: 0.875rem;">© 2025 All Rights Reserved</p>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()